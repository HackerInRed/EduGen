import os
from dotenv import load_dotenv
import streamlit as st
from PyPDF2 import PdfReader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.embeddings import FastEmbedEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
import asyncio
import json
import time
from docx import Document
from pptx import Presentation

# Load environment variables
load_dotenv()

def pdf_text(files):
    text = ""
    for file in files:
        reader = PdfReader(file)
        for page in reader.pages:
            text += page.extract_text()
    return text

def docx_text(files):
    text = ""
    for file in files:
        try:
            doc = Document(file)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
        except Exception as e:
            st.error(f"Error reading DOCX: {str(e)}")
    return text

def pptx_text(files):
    text = ""
    for file in files:
        try:
            ppt = Presentation(file)
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            st.error(f"Error reading PPTX: {str(e)}")
    return text

def txt_text(files):
    text = ""
    for file in files:
        try:
            text += str(file.read(), "utf-8") + "\n"
        except Exception as e:
            st.error(f"Error reading TXT: {str(e)}")
    return text

def process_files(files):
    text = ""
    for file in files:
        ext = file.name.split('.')[-1].lower()
        if ext == 'pdf':
            text += pdf_text([file])
        elif ext == 'docx':
            text += docx_text([file])
        elif ext == 'pptx':
            text += pptx_text([file])
        elif ext == 'txt':
            text += txt_text([file])
        else:
            st.warning(f"Unsupported file: {file.name}")
    return text

def chunk_text(text):
    splitter = RecursiveCharacterTextSplitter(chunk_size=50000, chunk_overlap=1000)
    chunks = splitter.split_text(text)
    return chunks

def create_store(chunks):
    embeddings = FastEmbedEmbeddings()
    store = FAISS.from_texts(chunks, embedding=embeddings)
    store.save_local("edugen")

def get_chain():
    template = """
    You are EduGen, an AI educational assistant.
    Provide clear answers based on the lecture materials.
    Context: {context}
    Question: {question}
    Answer:
    """
    
    try:
        asyncio.get_running_loop()
    except RuntimeError:
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
    
    model = ChatGoogleGenerativeAI(
        model="gemini-2.5-flash-lite",
        temperature=0.3
    )
    
    prompt = PromptTemplate(template=template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

def notes_chain():
    template = """
    Generate comprehensive notes from the lecture content.
    Include main topics, key concepts, and important points.
    Context: {context}
    Generate Notes:
    """
    
    try:
        asyncio.get_running_loop()
    except RuntimeError:
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
    
    model = ChatGoogleGenerativeAI(
        model="gemini-2.5-flash-lite",
        temperature=0.2
    )
    
    prompt = PromptTemplate(template=template, input_variables=["context"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

def quiz_chain():
    template = """
    Generate 5 multiple choice questions based on the content.
    Format as JSON array:
    [
        {{
            "question": "Question text?",
            "options": ["A", "B", "C", "D"],
            "correct_answer": 0,
            "explanation": "Why this is correct"
        }}
    ]
    Context: {context}
    Generate Quiz JSON:
    """
    
    try:
        asyncio.get_running_loop()
    except RuntimeError:
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
    
    model = ChatGoogleGenerativeAI(
        model="gemini-2.5-flash-lite",
        temperature=0.3
    )
    
    prompt = PromptTemplate(template=template, input_variables=["context"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

def cards_chain():
    template = """
    Generate 10 flashcards based on the content.
    Format as JSON array:
    [
        {{
            "front": "Term or concept",
            "back": "Definition or explanation"
        }}
    ]
    Context: {context}
    Generate Flashcards JSON:
    """
    
    try:
        asyncio.get_running_loop()
    except RuntimeError:
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
    
    model = ChatGoogleGenerativeAI(
        model="gemini-2.5-flash-lite",
        temperature=0.3
    )
    
    prompt = PromptTemplate(template=template, input_variables=["context"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

def get_docs(query="lecture content"):
    embeddings = FastEmbedEmbeddings()
    
    if os.path.exists("edugen"):
        db = FAISS.load_local("edugen", embeddings, allow_dangerous_deserialization=True)
        docs = db.similarity_search(query, k=10)
        return docs
    else:
        return None

def answer(question):
    docs = get_docs(question)
    if docs is None:
        return "Please upload documents first!"
    
    chain = get_chain()
    response = chain({
        "input_documents": docs,
        "question": question
    }, return_only_outputs=True)
    
    return response["output_text"]

def notes():
    docs = get_docs("comprehensive lecture summary")
    if docs is None:
        return "Please upload documents first!"
    
    chain = notes_chain()
    response = chain({
        "input_documents": docs,
        "question": "Generate notes"
    }, return_only_outputs=True)
    
    return response["output_text"]

def quiz():
    docs = get_docs("quiz questions")
    if docs is None:
        return None
    
    chain = quiz_chain()
    response = chain({
        "input_documents": docs,
        "question": "Generate quiz"
    }, return_only_outputs=True)
    
    try:
        text = response["output_text"]
        start = text.find('[')
        end = text.rfind(']') + 1
        if start != -1 and end != 0:
            json_str = text[start:end]
            data = json.loads(json_str)
            return data
        else:
            return None
    except json.JSONDecodeError:
        return None

def cards():
    docs = get_docs("key terms concepts")
    if docs is None:
        return None
    
    chain = cards_chain()
    response = chain({
        "input_documents": docs,
        "question": "Generate flashcards"
    }, return_only_outputs=True)
    
    try:
        text = response["output_text"]
        start = text.find('[')
        end = text.rfind(']') + 1
        if start != -1 and end != 0:
            json_str = text[start:end]
            data = json.loads(json_str)
            return data
        else:
            return None
    except json.JSONDecodeError:
        return None

def main():
    st.set_page_config(page_title="EduGen", page_icon="🎓", layout="wide", menu_items=None)

    st.markdown("""
    <style>
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    header {visibility: hidden;}
    .main .block-container {
        max-width: 100%;
        padding-top: 0.1rem;
        padding-bottom: 1rem;
        padding-left: 1rem;
        padding-right: 1rem;
        margin-top: 0.1rem;
    }
    .report-container {
        padding: 1rem;
        border-radius: 0.5rem;
        margin-top: 1rem;
    }
    button[title="View fullscreen"]{
        visibility: hidden;
    }
    </style>
    """, unsafe_allow_html=True)
    
    # Navigation buttons in header
    col1, col2, col3, col4, col5 = st.columns([1, 1, 1, 1, 1])
    
    upload_btn = col1.button("Upload Documents", use_container_width=True)
    notes_btn = col2.button("Summary Notes", use_container_width=True)
    quiz_btn = col3.button("Interactive Quiz", use_container_width=True)
    cards_btn = col4.button("Flashcards", use_container_width=True)
    qa_btn = col5.button("Q&A Assistant", use_container_width=True)

    # Initialize session state
    if "page" not in st.session_state:
        st.session_state.page = "upload"
    if "quiz_data" not in st.session_state:
        st.session_state.quiz_data = None
    if "quiz_started" not in st.session_state:
        st.session_state.quiz_started = False
    if "current_q" not in st.session_state:
        st.session_state.current_q = 0
    if "answers" not in st.session_state:
        st.session_state.answers = []
    if "score" not in st.session_state:
        st.session_state.score = 0
    if "completed" not in st.session_state:
        st.session_state.completed = False
    if "cards_data" not in st.session_state:
        st.session_state.cards_data = None
    if "current_card" not in st.session_state:
        st.session_state.current_card = 0
    if "show_back" not in st.session_state:
        st.session_state.show_back = False
    if "messages" not in st.session_state:
        st.session_state.messages = [
            {"role": "assistant", "content": "Hi! I'm EduGen. Upload your materials and ask me anything!"}
        ]

    # Handle navigation
    if upload_btn:
        st.session_state.page = "upload"
    elif notes_btn:
        st.session_state.page = "notes"
    elif quiz_btn:
        st.session_state.page = "quiz"
    elif cards_btn:
        st.session_state.page = "cards"
    elif qa_btn:
        st.session_state.page = "qa"

    # Upload page
    if st.session_state.page == "upload":
        st.markdown("<br>" * 3, unsafe_allow_html=True)
        st.markdown("<h2 style='text-align: center;'>Upload Your Lecture Documents</h2>", unsafe_allow_html=True)
        
        col1, col2, col3 = st.columns([1, 2, 1])
        with col2:
            files = st.file_uploader("Upload your lecture documents", accept_multiple_files=True, type=['pdf', 'docx', 'pptx', 'txt'])

        if files:
            with st.spinner("Processing files..."):
                text = process_files(files)
                if text.strip():
                    chunks = chunk_text(text)
                    create_store(chunks)
                    col1, col2, col3 = st.columns([1, 2, 1])
                    with col2:
                        st.success(f"{len(files)} files uploaded successfully!")
                else:
                    st.error("No text extracted from files")

    # Notes page
    elif st.session_state.page == "notes":
        if not os.path.exists("edugen"):
            st.warning("Please upload documents first.")
        else:
            st.header("Generate Summary Notes")
            st.markdown("Get comprehensive notes from your materials.")
            
            if st.button("Generate Notes", type="primary"):
                with st.spinner("Generating notes..."):
                    result = notes()
                    if result != "Please upload documents first!":
                        st.markdown("### Generated Notes")
                        st.markdown(f"<div class='report-container'>{result}</div>", unsafe_allow_html=True)
                        st.download_button(
                            label="Download Notes",
                            data=result,
                            file_name="notes.txt",
                            mime="text/plain"
                        )
                    else:
                        st.error(result)

    # Quiz page
    elif st.session_state.page == "quiz":
        if not os.path.exists("edugen"):
            st.warning("Please upload documents first.")
        else:
            st.header("Interactive Quiz")
            st.markdown("Test your knowledge with a quiz.")
            
            # Start quiz
            if not st.session_state.quiz_started or st.session_state.completed:
                if st.button("Start New Quiz", type="primary"):
                    with st.spinner("Generating quiz..."):
                        quiz_data = quiz()
                        if quiz_data is not None:
                            st.session_state.quiz_data = quiz_data
                            st.session_state.quiz_started = True
                            st.session_state.current_q = 0
                            st.session_state.answers = []
                            st.session_state.score = 0
                            st.session_state.completed = False
                            st.rerun()
                        else:
                            st.error("Failed to generate quiz. Try again!")
            
            # Quiz in progress
            elif st.session_state.quiz_started and not st.session_state.completed:
                quiz_data = st.session_state.quiz_data
                current = st.session_state.current_q
                
                if current < len(quiz_data):
                    question = quiz_data[current]
                    
                    st.progress(current / len(quiz_data))
                    st.markdown(f"**Question {current + 1} of {len(quiz_data)}**")
                    st.markdown(f"### {question['question']}")
                    
                    selected = st.radio(
                        "Choose your answer:",
                        options=range(len(question['options'])),
                        format_func=lambda x: f"{chr(65+x)}) {question['options'][x]}",
                        key=f"q_{current}"
                    )
                    
                    if st.button("Submit Answer", key=f"submit_{current}"):
                        st.session_state.answers.append(selected)
                        
                        if selected == question['correct_answer']:
                            st.session_state.score += 1
                            st.success("Correct!")
                        else:
                            correct_option = question['options'][question['correct_answer']]
                            st.error(f"Incorrect. Correct answer: {chr(65+question['correct_answer'])}) {correct_option}")
                        
                        st.info(f"**Explanation:** {question['explanation']}")
                        st.session_state.current_q += 1
                        
                        if st.session_state.current_q >= len(quiz_data):
                            st.session_state.completed = True
                        
                        time.sleep(2)
                        st.rerun()
            
            # Quiz completed
            if st.session_state.completed:
                st.markdown("## Quiz Completed!")
                score = st.session_state.score
                total = len(st.session_state.quiz_data)
                percentage = (score / total) * 100
                
                st.metric("Your Score", f"{score}/{total} ({percentage:.1f}%)")
                
                if percentage >= 90:
                    st.success("Excellent! Great understanding!")
                elif percentage >= 70:
                    st.success("Good job! You understand most concepts.")
                elif percentage >= 50:
                    st.warning("Not bad, but review the material.")
                else:
                    st.error("Study the material more thoroughly.")

    # Flashcards page
    elif st.session_state.page == "cards":
        if not os.path.exists("edugen"):
            st.warning("Please upload documents first.")
        else:
            st.header("Interactive Flashcards")
            st.markdown("Study key concepts with flashcards.")
            
            if st.session_state.cards_data is None:
                if st.button("Generate Flashcards", type="primary"):
                    with st.spinner("Generating flashcards..."):
                        cards_data = cards()
                        if cards_data is not None:
                            st.session_state.cards_data = cards_data
                            st.session_state.current_card = 0
                            st.session_state.show_back = False
                            st.rerun()
                        else:
                            st.error("Failed to generate flashcards!")
            
            if st.session_state.cards_data is not None:
                cards_data = st.session_state.cards_data
                current = st.session_state.current_card
                
                if current < len(cards_data):
                    st.progress((current + 1) / len(cards_data))
                    st.markdown(f"**Flashcard {current + 1} of {len(cards_data)}**")
                    
                    card = cards_data[current]
                    
                    card_style = """
                    <div style="
                        background: linear-gradient(145deg, #ffffff, #f8f9fa);
                        padding: 2.5rem;
                        border-radius: 20px;
                        box-shadow: 0 8px 32px rgba(0, 0, 0, 0.1);
                        text-align: center;
                        margin: 1.5rem 0;
                        min-height: 250px;
                        display: flex;
                        align-items: center;
                        justify-content: center;
                    ">
                        <div style="
                            font-size: 1.3em; 
                            font-weight: 600;
                            color: #2c3e50;
                            line-height: 1.6;
                        ">
                            {}
                        </div>
                    </div>
                    """
                    
                    if not st.session_state.show_back:
                        st.markdown(card_style.format(card['front']), unsafe_allow_html=True)
                        col1, col2, col3 = st.columns([1, 2, 1])
                        with col2:
                            if st.button("Flip to Answer", type="primary", use_container_width=True):
                                st.session_state.show_back = True
                                st.rerun()
                    else:
                        st.markdown(card_style.format(card['back']), unsafe_allow_html=True)
                        col1, col2, col3 = st.columns([1, 1, 1])
                        with col1:
                            if st.button("Flip Back", use_container_width=True):
                                st.session_state.show_back = False
                                st.rerun()
                        with col2:
                            if st.button("Need Review", use_container_width=True):
                                st.session_state.show_back = False
                                st.info("Marked for review!")
                                st.rerun()
                        with col3:
                            if st.button("Got It!", use_container_width=True):
                                st.session_state.show_back = False
                                st.success("Great job!")
                                st.rerun()
                    
                    # Navigation
                    st.markdown("---")
                    col1, col2, col3, col4, col5 = st.columns([1, 1, 1, 1, 1])
                    
                    with col1:
                        if current > 0:
                            if st.button("← Previous", use_container_width=True):
                                st.session_state.current_card -= 1
                                st.session_state.show_back = False
                                st.rerun()
                    
                    with col3:
                        if st.button("Start Over", use_container_width=True):
                            st.session_state.current_card = 0
                            st.session_state.show_back = False
                            st.rerun()
                    
                    with col4:
                        if st.button("New Set", use_container_width=True):
                            st.session_state.cards_data = None
                            st.session_state.current_card = 0
                            st.session_state.show_back = False
                            st.rerun()
                    
                    with col5:
                        if current < len(cards_data) - 1:
                            if st.button("Next →", use_container_width=True):
                                st.session_state.current_card += 1
                                st.session_state.show_back = False
                                st.rerun()
                
                else:
                    st.success("You've completed all flashcards!")
                    if st.button("Study Again"):
                        st.session_state.current_card = 0
                        st.session_state.show_back = False
                        st.rerun()

    # Q&A page
    elif st.session_state.page == "qa":
        if not os.path.exists("edugen"):
            st.warning("Please upload documents first.")
        else:
            st.header("Q&A Assistant")
            st.markdown("Ask questions about your materials.")
            
            for message in st.session_state.messages:
                with st.chat_message(message["role"]):
                    st.write(message["content"])
            
            if prompt := st.chat_input("Ask a question..."):
                st.session_state.messages.append({"role": "user", "content": prompt})
                
                with st.chat_message("user"):
                    st.write(prompt)
                
                with st.chat_message("assistant"):
                    with st.spinner("Thinking..."):
                        response = answer(prompt)
                        st.write(response)
                        st.session_state.messages.append({"role": "assistant", "content": response})
            
            if st.button("Clear Chat"):
                st.session_state.messages = [
                    {"role": "assistant", "content": "Hi! I'm EduGen. Upload your materials and ask me anything!"}
                ]
                st.rerun()

if __name__ == "__main__":
    main()
